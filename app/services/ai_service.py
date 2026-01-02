# app/services/ai_service.py

import json
import re
import logging
from io import BytesIO
from app.utils import ai_utils
from groq import Groq
import litellm
from litellm import completion

litellm.drop_params = True

from app.utils.ai_utils import get_smart_model, clean_json_output



# Setup Logger
logger = logging.getLogger(__name__)
llm_client = Groq()



# Coba import pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

except litellm.RateLimitError:
    logger.error("Quota habis. Silakan upgrade paket atau hubungi admin.")

    

class AIService:
    
    # ==========================================
    # 0. HELPER: SMART CONTEXT PRUNING
    # ==========================================
    @staticmethod
    def _prune_text(text, max_chars=10000, position='end'):
        """
        Memotong teks agar tidak meledakkan token limit AI.
        """
        if not text:
            return ""
        
        if len(text) <= max_chars:
            return text
            
        if position == 'end':
            return "... (context truncated) ...\n" + text[-max_chars:]
        else:
            return text[:max_chars] + "\n... (references truncated) ..."

    # =========================================================================
    # 1. CORE: ROBUST PROMPT ENGINEERING (OTAK BARU - METHODOLOGY AWARE)
    # =========================================================================
    @staticmethod
    def _construct_robust_prompt(task_type, data):
        """
        Merakit System Prompt & User Prompt.
        Disini aturan main (Rules) diterapkan.
        """
        # Ambil data dasar dengan nilai default aman
        user_input = data.get('content') or data.get('prompt') or data.get('message', '')
        context_str = data.get('context', '') or data.get('context_material', '')
        
        # Ambil konteks spesifik project
        p_title = data.get('context_title', 'Topik Umum')
        p_problem = data.get('context_problem', '-')
        
        # ==========================================
        # 1. LOGIKA CHAT (Limit 200 Kata)
        # ==========================================
        if task_type == 'chat':
            system_instruction = (
                "Anda adalah Asisten Riset Akademik (Writing Studio AI).\n"
                "PERATURAN UTAMA (STRICT):\n"
                "1. Jawablah pertanyaan pengguna dengan gaya bahasa akademis, formal, namun mudah dimengerti.\n"
                "2. BATAS PANJANG: Jawaban Anda TIDAK BOLEH lebih dari 200 kata per balasan.\n"
                "3. Jika penjelasan membutuhkan ruang lebih panjang, sarankan pengguna menggunakan fitur 'Generator' (Tools Tab).\n"
                "4. Gunakan Bahasa Indonesia baku (EYD).\n"
                "5. Jika ada konteks proyek di bawah, gunakan itu sebagai referensi jawaban."
            )
            
            if context_str:
                system_instruction += f"\n\n[KONTEKS DOKUMEN/PROYEK]:\n{context_str}"
            
            if p_title:
                system_instruction += f"\n\n[INFO SKRIPSI USER]:\nJudul: {p_title}\nMasalah: {p_problem}"
                
            return {'system': system_instruction, 'user': user_input}

        # ==========================================
        # 2. LOGIKA MIKRO (Paraphrase, Expand)
        # ==========================================
        elif task_type in ['paraphrase', 'expand_text', 'formalize_text']:
            style = data.get('style', 'academic')
            modes = {
                'paraphrase': "Tulis ulang teks berikut agar lolos plagiasi (Turnitin) namun makna tetap sama.",
                'expand_text': "Kembangkan kalimat berikut menjadi paragraf yang lebih detail dan berbobot.",
                'formalize_text': "Ubah teks berikut menjadi bahasa akademis formal (baku)."
            }
            instruction = modes.get(task_type, "Perbaiki teks berikut.")
            
            system_instruction = (
                f"Tugas: {instruction}\n"
                f"Gaya Bahasa: {style}.\n"
                "Aturan: Langsung berikan hasil revisi tanpa kata pengantar 'Berikut adalah...'."
            )
            return {'system': system_instruction, 'user': user_input}

        # ==========================================
        # 3. LOGIKA BERAT (Bab 1 - 5)
        # ==========================================
        # Variabel bantu untuk logic prompt skripsi
        # Kita definisikan di sini agar aman dipakai di blok bawah
        topic = data.get('topic') or p_title or 'Topik Umum'
        method_name = data.get('methodology') or data.get('context_method') or 'Metode Standar'
        ref_str = f"Referensi: {context_str[:1000]}..." if context_str else ""
        problem = data.get('problem') or p_problem or 'Belum ditentukan'
        
        # Cek tipe metodologi (untuk percabangan prompt)
        is_kualitatif = 'kualitatif' in str(method_name).lower() or 'qualitative' in str(method_name).lower()

        base_academic_system = (
            "Anda adalah Dosen Pembimbing Senior & Pakar Metodologi Penelitian.\n"
            "Tugas Anda adalah membantu mahasiswa menulis draf skripsi bagian per bagian.\n"
            "Gaya Penulisan: Ilmiah, Obyektif, Analitis, dan Koheren.\n"
            "Gunakan referensi jika diberikan dalam konteks."
        )

        # ---------------------------------------------------------------------
        # A. BAB 1: PENDAHULUAN
        # ---------------------------------------------------------------------
        if task_type == 'bab1_latar_belakang' or task_type == 'background':
            return {
                "system": "Anda adalah Dosen Metodologi Penelitian Senior yang kritis.",
                "user": f"""
                TUGAS: Tulis Latar Belakang Masalah (Sekitar 400 kata).
                JUDUL: "{topic}"
                METODE: {method_name}
                MASALAH UTAMA: "{problem}"
                
                STRUKTUR SEGITIGA TERBALIK:
                1. Das Sollen (Harapan/Ideal): Kondisi ideal fenomena ini secara global/umum.
                2. Das Sein (Fakta/Masalah): Data kesenjangan atau masalah riil di lapangan (Gap).
                3. Urgensi: Mengapa masalah ini harus diteliti sekarang?
                
                STRICT RULES:
                - Hindari kalimat klise "Pada zaman globalisasi...". Langsung ke topik.
                - Sertakan data faktual jika memungkinkan (buat placeholder [DATA] jika tidak tau).
                - Bahasa harus baku, objektif, dan denotatif.
                
                {ref_str}
                """
            }
        
        elif task_type == 'problem_formulation':
            return {
                "system": "Anda ahli merumuskan pertanyaan penelitian.",
                "user": f"""
                JUDUL: "{topic}"
                METODE: {method_name}
                TUGAS: Buat Rumusan Masalah dalam bentuk poin-poin pertanyaan.
                Pastikan relevan dengan judul dan variabel penelitian.
                """
            }

        # ---------------------------------------------------------------------
        # B. BAB 2: TINJAUAN PUSTAKA (Spesifik per Sub-Bab)
        # ---------------------------------------------------------------------
        elif task_type == 'content_specific': 
            sub_bab = data.get('sub_bab', 'Sub Bab')
            points = data.get('points', [])
            points_str = "\n".join([f"- {p}" for p in points])
            
            # DETEKSI HIPOTESIS (Strict Rule Kuantitatif Only)
            is_hipotesis = "hipotesis" in sub_bab.lower()
            
            prompt_tambahan = ""
            if is_hipotesis:
                if is_kualitatif:
                    return {
                        "system": "System Alert.",
                        "user": "Tulis: 'Dalam penelitian Kualitatif, tidak digunakan Hipotesis statistik, melainkan Pertanyaan Penelitian atau Proposisi. Bagian ini tidak relevan.'"
                    }
                else:
                    prompt_tambahan = """
                    ATURAN KHUSUS HIPOTESIS (KUANTITATIF):
                    - JANGAN bertele-tele. Langsung tulis rumusan hipotesis.
                    - Gunakan format baku:
                      H0: Tidak terdapat pengaruh...
                      H1: Terdapat pengaruh...
                    - Jelaskan dasar dugaan (premis) singkat sebelum merumuskan.
                    """

            return {
                "system": "Anda adalah Penulis Akademik Spesialis Literatur Review.",
                "user": f"""
                KONTEKS SKRIPSI:
                Judul: "{topic}"
                Metode: {method_name}
                
                TUGAS: Tulis konten untuk Sub-Bab: "{sub_bab}".
                POIN:
                {points_str}
                
                {ref_str}
                {prompt_tambahan}

                ATURAN PENULISAN:
                1. Fokus HANYA pada sub-bab ini.
                2. Sintesiskan teori (jangan cuma list definisi).
                3. WAJIB SITASI (Author, Year) dari daftar referensi di atas jika relevan.
                4. Panjang: 300-500 kata.
                """
            }

        # ---------------------------------------------------------------------
        # C. BAB 3: METODOLOGI (Sesuai Jenis Penelitian)
        # ---------------------------------------------------------------------
        elif task_type == 'bab3_metode' or task_type == 'method_type':
            if is_kualitatif:
                # PROMPT KUALITATIF
                instruction = """
                FOKUS KUALITATIF:
                1. Pendekatan: Deskriptif Kualitatif / Fenomenologi / Studi Kasus (Pilih yang cocok dengan judul).
                2. Subjek: Gunakan istilah 'Informan' atau 'Narasumber' (BUKAN Responden).
                3. Teknik Sampling: Purposive Sampling atau Snowball Sampling.
                4. Keabsahan Data: Triangulasi (Sumber/Teknik/Waktu).
                5. Analisis: Reduksi data, Penyajian data, Penarikan kesimpulan (Miles & Huberman).
                """
            else:
                # PROMPT KUANTITATIF
                instruction = """
                FOKUS KUANTITATIF:
                1. Pendekatan: Asosiatif / Komparatif / Eksperimen.
                2. Objek: Gunakan istilah 'Populasi' dan 'Sampel' (Responden).
                3. Teknik Sampling: Random Sampling / Stratified / dll.
                4. Instrumen: Kuesioner/Angket (Uji Validitas & Reliabilitas).
                5. Analisis: Statistik Deskriptif & Inferensial (Uji T / Uji F / Regresi).
                """

            return {
                "system": f"Anda adalah Ahli Metodologi Penelitian {method_name}.",
                "user": f"""
                TUGAS: Tulis Bab 3 Metodologi Penelitian yang presisi.
                JUDUL: "{topic}"
                
                {instruction}
                
                ATURAN FATAL:
                - JANGAN tertukar istilah Kualitatif vs Kuantitatif.
                - Jelaskan alasan pemilihan metode tersebut untuk judul ini.
                """
            }

        # ---------------------------------------------------------------------
        # D. BAB 4: HASIL & PEMBAHASAN (Interpretasi Data)
        # ---------------------------------------------------------------------
        elif task_type == 'interpret_data':
            raw_data = json.dumps(data.get('raw_data', {}), indent=2)
            
            if is_kualitatif:
                # LOGIC KUALITATIF (Interpretasi Teks/Tema)
                return {
                    "system": "Anda adalah Peneliti Kualitatif.",
                    "user": f"""
                    TUGAS: Buat Pembahasan Bab 4 Kualitatif.
                    JUDUL: "{topic}"
                    
                    DATA TEMUAN (Tema/Kutipan):
                    {raw_data}
                    
                    INSTRUKSI:
                    1. Lakukan 'Thick Description' (Deskripsi Mendalam).
                    2. Hubungkan temuan lapangan dengan Teori di Bab 2.
                    3. Jangan bicara angka statistik/persentase. Bicara pola/makna.
                    """
                }
            else:
                # LOGIC KUANTITATIF (Interpretasi Angka)
                return {
                    "system": "Anda adalah Data Analyst Statistik.",
                    "user": f"""
                    TUGAS: Interpretasi Data Statistik Bab 4.
                    JUDUL: "{topic}"
                    Hipotesis: "{data.get('context_hypothesis', 'H1')}"
                    
                    DATA STATISTIK (JSON):
                    {raw_data}
                    
                    INSTRUKSI FATAL:
                    1. BACA ANGKA DENGAN TEPAT (Signifikansi, t-hitung, R-square).
                    2. JANGAN MENGARANG ANGKA! Gunakan hanya yang ada di data.
                    3. Tentukan Hipotesis Diterima/Ditolak berdasarkan data (misal: sig < 0.05).
                    4. Bahas implikasi hasil ini terhadap variabel penelitian.
                    """
                }

        # ---------------------------------------------------------------------
        # E. BAB 5: PENUTUP
        # ---------------------------------------------------------------------
        elif task_type == 'bab5_penutup':
            return {
                "system": "Anda adalah Penulis Penutup Skripsi.",
                "user": f"""
                TUGAS: Tulis Kesimpulan dan Saran.
                JUDUL: "{topic}"
                METODE: {method_name}
                
                INSTRUKSI:
                1. Kesimpulan: Jawab rumusan masalah secara singkat. { "Jangan pakai angka statistik lagi." if is_kualitatif else "Sebutkan hasil hipotesis secara ringkas." }
                2. Saran: Berikan rekomendasi operasional untuk objek penelitian.
                """
            }

        # ---------------------------------------------------------------------
        # FALLBACK DEFAULT (Untuk task umum/lainnya)
        # ---------------------------------------------------------------------
        # FIX: Variabel 'topic', 'method_name', 'ref_str' sudah didefinisikan di awal fungsi (Line ~90)
        # Jadi aman dipakai di sini.
        return {
            "system": "Anda asisten penulisan skripsi.",
            "user": f"Bantu saya menulis tentang: {topic}. Metode: {method_name}. {ref_str}. Detail Instruksi: {user_input}"
        }

    @staticmethod
    def writing_assistant_stream(user, data):
        """
        Fungsi Utama Generator (Streaming).
        """
        input_data = data.get('data', {})
        task_type = data.get('task', 'general')

        # 1. PILIH MODEL (The Magic Switch)
        selected_model = get_smart_model(task_type, user)
        
        logger.info(f"Generating for {user.id} | Task: {task_type} | Model: {selected_model}")

        # 2. RAKIT PROMPT
        prompt_pack = AIService._construct_robust_prompt(task_type, input_data)

        try:
            # 3. PANGGIL API (LiteLLM)
            response = completion(
                model=selected_model,
                messages=[
                    {"role": "system", "content": prompt_pack['system']},
                    {"role": "user", "content": prompt_pack['user']}
                ],
                temperature=0.3, # Rendah agar stabil
                stream=True,     # Wajib true untuk efek mengetik
                max_tokens=2500,  # Buffer cukup panjang
                drop_params=True # <--- FIX: Drop params yang tidak didukung (seperti temp 0.3 di model tertentu)
            )
            
            # 4. STREAM HASIL
            for chunk in response:
                content = chunk.choices[0].delta.content
                if content:
                    yield content

        except Exception as e:
            logger.error(f"AI Stream Error: {e}")
            yield f"<span style='color:red'>[Sistem Error]: {str(e)}</span>"

    # =========================================================================
    # 3. GENERATE OUTLINE BAB 2 (UPDATED: METHODOLOGY AWARE)
    # =========================================================================
    @staticmethod
    def generate_bab2_outline(user, payload):
        judul = payload.get("judul_penelitian")
        method = payload.get("methodology", "quantitative").lower() 
        
        if not judul: raise ValueError("Judul wajib diisi")

        bagian_akhir = '"2.4 Fokus Penelitian / Proposisi": ["<uraian>"]' if "kualitatif" in method else '"2.4 Hipotesis": ["<uraian>"]'

        # STRICT PROMPT untuk mencegah model "curhat" atau bikin draf
        prompt = f"""
        TUGAS: Buat STRUKTUR OUTLINE BAB 2 (Tinjauan Pustaka) untuk judul: "{judul}".
        METODE: {method}.

        ATURAN SANGAT PENTING (JIKA MELANGGAR = ERROR):
        1. JANGAN MENULIS ISI/DRAF SKRIPSI. HANYA STRUKTUR (JUDUL SUB-BAB & POIN).
        2. KELUARKAN HANYA JSON. TIDAK BOLEH ADA TEKS PEMBUKA/PENUTUP.
        3. Format harus persis seperti contoh di bawah.

        CONTOH FORMAT JSON (WAJIB DIIKUTI):
        {{
            "outline": {{
                "2.1 Konsep Variabel X": ["Definisi menurut ahli", "Indikator", "Faktor yang mempengaruhi"],
                "2.2 Konsep Variabel Y": ["Teori pendukung", "Dimensi", "Dampak"],
                "2.3 Penelitian Terdahulu": ["Matrix penelitian relevan", "Perbedaan dengan penelitian ini"],
                "2.4 Kerangka Berpikir": ["Hubungan antar variabel"],
                {bagian_akhir}
            }}
        }}
        """
        
        try:
            # Gunakan model yang sedikit lebih pintar atau mode JSON strict jika didukung
            response = completion(
                model="groq/llama-3.1-70b-versatile", # Upgrade model dikit biar lebih patuh
                messages=[
                    {"role": "system", "content": "You are a JSON generator API. You ONLY output valid JSON. No conversational text."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1, # Sangat rendah agar deterministik
                response_format={"type": "json_object"},
                drop_params=True,
                num_retries=2 
            )
            
            content = response.choices[0].message.content
            
            # --- FIX: Extra Cleaning untuk membuang teks sampah ---
            # Cari kurung kurawal pertama '{' dan terakhir '}'
            match = re.search(r'\{.*\}', content, re.DOTALL)
            if match:
                clean_json = match.group(0)
            else:
                clean_json = clean_json_output(content)
            
            parsed_data = json.loads(clean_json)
            
            # Normalisasi output menjadi array standar
            final_outline_array = []
            outline_data = parsed_data.get("outline", {})
            
            # Jika model malah mengembalikan array langsung (kadang terjadi)
            if isinstance(outline_data, list):
                return outline_data

            for key, val in outline_data.items():
                final_outline_array.append({
                    "sub_bab": key, 
                    "poin_pembahasan": val if isinstance(val, list) else [str(val)]
                })
            
            return final_outline_array

        except Exception as e:
            logger.error(f"Generate Outline Fatal Error: {e}. Raw Content: {content if 'content' in locals() else 'None'}")
            # Fallback agar tidak crash
            return [
                {"sub_bab": "2.1 Landasan Teori", "poin_pembahasan": ["Definisi Variabel Utama", "Teori Pendukung"]},
                {"sub_bab": "2.2 Penelitian Terdahulu", "poin_pembahasan": ["Kajian empiris relevan"]},
                {"sub_bab": "2.3 Kerangka Pemikiran", "poin_pembahasan": ["Alur logika penelitian"]}
            ]

    # ==========================================
    # 4. CONTEXT AUTOPILOT
    # ==========================================
    @staticmethod
    def generate_specific_context_field(field_type, context):
        title = context.get('title', 'Tanpa Judul')
        loc = context.get('research_location', 'Lokasi Penelitian')
        
        prompt_guide = "Buatkan ringkasan singkat."
        if field_type == 'problem_statement':
            prompt_guide = "Buatkan 3 poin Rumusan Masalah kritis."
        elif field_type == 'research_objectives':
            prompt_guide = "Buatkan 3 poin Tujuan Penelitian."
        elif field_type == 'significance':
            prompt_guide = "Buatkan Manfaat Teoretis & Praktis."
        elif field_type == 'variables':
            prompt_guide = "Tentukan Variabel X, Y, dan definisinya."
        elif field_type == 'hypothesis':
            prompt_guide = "Rumuskan Hipotesis (H1, H2) yang logis."
        elif field_type == 'methodology':
            prompt_guide = "Sarankan Metode (Pendekatan, Desain, Teknik)."
        elif field_type == 'population_sample':
            prompt_guide = f"Sarankan Populasi & Sampel untuk lokasi {loc}."
        elif field_type == 'data_analysis':
            prompt_guide = "Sarankan Teknik Analisis Data."

        system_instruction = (
            f"Anda Konsultan Riset. Judul: {title}. Lokasi: {loc}.\n"
            "Output: Poin-poin akademis, langsung, baku."
        )

        try:
            response = completion(
                model="groq/llama-3.1-8b-instant",
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": f"Field: {field_type}\n{prompt_guide}"}
                ],
                temperature=0.3
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error: {str(e)}"

    # ==========================================
    # 5. MAGIC EXPAND & PARAPHRASE
    # ==========================================
    @staticmethod
    def expand_stream(user, text, context_references=""):
        ref_instruction = ""
        if context_references:
            ref_instruction = f"\n[REFERENSI]:\n{context_references}\nGUNAKAN REFERENSI INI."
        else:
            ref_instruction = "\nGunakan pengetahuan akademis umum & sitasi umum jika perlu."

        system_instruction = (
            "Anda asisten skripsi. TUGAS: Kembangkan kalimat user jadi 1 PARAGRAF padat (Max 150 kata).\n"
            "WAJIB ADA SITASI (Bodynote) format (Nama, Tahun).\n"
            f"{ref_instruction}\n"
            "BAHASA: Indonesia Baku."
        )

        try:
            response = completion(
                model="groq/llama-3.1-8b-instant",
                messages=[{"role": "system", "content": system_instruction}, {"role": "user", "content": text}],
                stream=True, temperature=0.3
            )
            for chunk in response:
                if chunk.choices[0].delta.content: yield chunk.choices[0].delta.content
        except Exception as e:
            yield f"[Error: {str(e)}]"

    @staticmethod
    def paraphrase_stream(user, text, style='academic', model='fast'):
        system_instruction = (
            "Anda editor jurnal. Parafrase teks berikut agar FORMAL, BAKU, dan EFEKTIF.\n"
            "JANGAN translate ke Inggris. Tetap Bahasa Indonesia.\n"
            "Pertahankan panjang tulisan."
        )
        try:
            response = completion(
                model="groq/llama-3.1-8b-instant", 
                messages=[{"role": "system", "content": system_instruction}, {"role": "user", "content": text}],
                stream=True, temperature=0.2
            )
            for chunk in response:
                if chunk.choices[0].delta.content: yield chunk.choices[0].delta.content
        except Exception as e:
            yield f"[Error: {str(e)}]"

    @staticmethod
    def chat_stream(user, message, project_id=None):
        return ai_utils.get_chat_response_stream(message)

    # ==========================================
    # 6. DOSEN KILLER / AUDIT & SIMULASI
    # ==========================================
    @staticmethod
    def ask_ai_json(user, prompt, context=None):
        system_instruction = "Bertindak sebagai Dosen Pembimbing. Output JSON: [{ 'target': '', 'type': 'critical', 'feedback': '', 'fix': '' }]"
        try:
            response = completion(
                model="groq/llama-3.3-70b-versatile",
                messages=[{"role": "system", "content": system_instruction}, {"role": "user", "content": prompt}],
                response_format={"type": "json_object"}
            )
            parsed = json.loads(response.choices[0].message.content)
            result = parsed.get('reviews') or parsed.get('items') or [parsed]
            return [i for i in result if 'target' in i]
        except: return []

    @staticmethod
    def generate_logic_matrix(user, problem_text, conclusion_text):
        system_instruction = """
        Bandingkan Rumusan Masalah vs Kesimpulan.
        Output JSON: { "summary": "...", "score": 0-100, "matrix": [{"poin_masalah": "...", "poin_kesimpulan": "...", "status": "MATCH/MISMATCH"}] }
        """
        try:
            response = completion(
                model="groq/llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": f"MASALAH: {problem_text}\nKESIMPULAN: {conclusion_text}"}
                ],
                response_format={"type": "json_object"}
            )
            return json.loads(response.choices[0].message.content)
        except Exception as e:
            return {"summary": str(e), "score": 0, "matrix": []}

    @staticmethod
    def thesis_defense_simulation(user, action, data):
        """
        Simulasi Sidang Skripsi Interaktif.
        Action: 'start', 'answer', 'evaluate'
        """
        # Gunakan model pintar untuk user Pro, atau standar untuk free
        model_name = get_smart_model('defense_simulation', user)
        
        # 1. START SESSION
        if action == 'start':
            ex_type = data.get('examiner_type', 'critical')
            difficulty = data.get('difficulty', 'hard')
            context = data.get('project_context', {})
            
            # --- TUNING DOSEN KILLER ---
            base_role = "Anda adalah Dosen Penguji Skripsi Senior."
            tone = ""
            
            if ex_type == 'critical':
                tone = """
                KARAKTER: SANGAT KRITIS, CYNICAL, DAN JUDES.
                ATURAN BICARA:
                1. JANGAN MEMUJI. JANGAN BASA-BASI.
                2. Langsung serang kelemahan logika judul/masalah.
                3. Pertanyaan harus PENDEK (Max 2 kalimat). Bikin mahasiswa deg-degan.
                4. Jika jawaban mahasiswa lemah, potong pembicaraan dengan kata seperti "Logikanya di mana?", "Anda yakin?".
                """
            elif ex_type == 'methodologist':
                tone = """
                KARAKTER: Ahli Metodologi yang kaku.
                ATURAN BICARA:
                1. Fokus hanya pada validitas data, rumus, dan sampling.
                2. Pertanyaan sangat teknis dan menjebak.
                3. Max 30 kata per pertanyaan.
                """
            else:
                tone = "KARAKTER: Supportif tapi tegas. Menguji pemahaman konsep dasar."

            system_instruction = f"{base_role}\n{tone}\n\nKonteks Skripsi Mahasiswa:\nJudul: {context.get('title')}\nMasalah: {context.get('problem')}\nMetode: {context.get('method')}"

            prompt = f"Sidang dimulai. Mahasiswa baru masuk ruangan. Tembak pertanyaan pertama tentang urgensi atau masalah penelitian. JANGAN SAPA 'HALO'."
            
            try:
                response = completion(
                    model=model_name,
                    messages=[{"role": "system", "content": system_instruction}, {"role": "user", "content": prompt}],
                    temperature=0.8, # Tinggi biar variatif dan galak
                    drop_params=True,
                    num_retries=2
                )
                return {"message": response.choices[0].message.content}
            except Exception as e:
                logger.error(f"Defense Start Error: {e}")
                return {"message": "Dosen sedang sibuk. (System Error)"}

        # 2. ANSWER & FOLLOW-UP
        elif action == 'answer':
            history = data.get('history', []) # List of {role, content}
            user_answer = data.get('answer', '')
            ex_type = data.get('examiner_type', 'critical')
            
            # Re-construct system prompt (karena stateless)
            system_role = "Anda Dosen Penguji Killer. Tugas: Cari celah jawaban mahasiswa. Pertanyaan harus singkat, pedas, dan tajam (Max 30 kata). Jangan pernah bilang 'Bagus'."
            if ex_type == 'methodologist': system_role = "Anda Dosen Metodologi. Kejar teknisnya. Pertanyaan singkat."
            
            # Format history untuk LLM
            messages = [{"role": "system", "content": system_role}]
            
            # Ambil max 6 chat terakhir biar hemat token & fokus konteks baru
            recent_history = history[-6:] 
            for msg in recent_history:
                messages.append({"role": "assistant" if msg['role'] == 'assistant' else "user", "content": msg['content']})
            
            messages.append({"role": "user", "content": user_answer})
            
            try:
                response = completion(
                    model=model_name,
                    messages=messages,
                    temperature=0.7,
                    drop_params=True,
                    num_retries=2
                )
                # Dosen tidak memberikan skor di chat, hanya teks
                return {"message": response.choices[0].message.content} 
            except Exception as e:
                return {"message": "..."}

        # 3. EVALUATE SESSION
        elif action == 'evaluate':
            history = data.get('history', [])
            
            transcript = "\n".join([f"{m['role'].upper()}: {m['content']}" for m in history])
            
            system_prompt = """
            PERAN: Ketua Sidang Skripsi.
            TUGAS: Berikan evaluasi akhir berdasarkan transkrip sidang.
            
            OUTPUT WAJIB JSON:
            {
                "score": (Angka 0-100),
                "verdict": "LULUS" / "REVISI TOTAL" / "TIDAK LULUS",
                "strengths": "Poin positif singkat...",
                "weaknesses": "Poin negatif singkat...",
                "advice": "Saran perbaikan..."
            }
            """
            
            try:
                response = completion(
                    model=model_name,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"TRANSKRIP:\n{transcript}"}
                    ],
                    response_format={"type": "json_object"},
                    drop_params=True,
                    num_retries=2
                )
                content = response.choices[0].message.content
                return json.loads(clean_json_output(content))
            except Exception as e:
                logger.error(f"Defense Eval Error: {e}")
                return {"score": 0, "verdict": "ERROR", "weaknesses": "Gagal evaluasi."}

    # ==========================================
    # 7. UTILS: PPT, STYLE, COMPLIANCE
    # ==========================================
    @staticmethod
    def generate_ppt(user, data):
        if not Presentation: raise Exception("Install python-pptx")
        prs = Presentation()
        # Slide Judul
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = data.get('title', 'Presentasi Skripsi')
        slide.placeholders[1].text = f"Oleh: {data.get('student_name', '')}"
        
        # Slide Isi
        for item in data.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title: slide.shapes.title.text = item.get('title', '')
            tf = slide.placeholders[1].text_frame
            for p in item.get('points', []): tf.add_paragraph().text = p
            
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    @staticmethod
    def check_method_compliance(text, method_mode):
        """Cek istilah salah kamar (e.g. Responden di Kualitatif)."""
        issues = []
        rules = {}
        if method_mode == 'quantitative':
            rules = {"informan": "responden", "narasumber": "responden", "makna": "pengaruh"}
        elif method_mode == 'qualitative':
            rules = {"responden": "informan", "hipotesis": "fokus penelitian", "pengaruh": "relevansi"}
        
        text_lower = text.lower()
        for bad, fix in rules.items():
            pattern = r'\b' + re.escape(bad) + r'\b'
            for m in re.finditer(pattern, text_lower):
                issues.append({
                    "target": text[m.start():m.end()],
                    "type": "critical",
                    "feedback": f"Istilah '{bad}' tidak cocok untuk metode {method_mode}.",
                    "fix": fix
                })
        return issues

    @staticmethod
    def analyze_logic_flow(data):
        """
        Menganalisis koherensi antar bab (Benang Merah).
        """
        context = data.get('context', {})
        title = context.get('title', '')
        problem = context.get('problem', '')
        objectives = context.get('objectives', '')
        method = context.get('method', '')
        conclusion = context.get('conclusion', '') 

        # Paksa pakai model validasi yang kuat
        model_name = ai_utils.get_optimal_model_for_task('logic_check')
        
        system_instruction = (
            "Anda adalah Auditor Logika Akademik. Tugas Anda mencari 'Logical Fallacy' atau inkonsistensi."
            "Output WAJIB format JSON."
        )

        user_prompt = f"""
        DATA SKRIPSI:
        1. Judul: "{title}"
        2. Rumusan Masalah: "{problem}"
        3. Tujuan Penelitian: "{objectives}"
        4. Metode: "{method}"
        5. Kesimpulan: "{conclusion or "(Belum ditulis)"}"

        TUGAS AUDIT:
        Lakukan Cross-Check berikut:
        A. (Judul vs Masalah): Apakah variabel di Judul konsisten dengan Masalah?
        B. (Masalah vs Tujuan): Apakah jumlah poin sinkron?
        C. (Metode vs Masalah): Apakah metodenya cocok?
        D. (Tujuan vs Kesimpulan): Apakah kesimpulan menjawab tujuan?

        OUTPUT JSON:
        {{
            "overall_score": (0-100),
            "verdict": "Solid / Perlu Revisi / Fatal",
            "summary": "Ringkasan penilaian.",
            "checks": [
                {{ "pair": "Judul vs Masalah", "status": "Valid" | "Warning" | "Critical", "feedback": "..." }},
                {{ "pair": "Masalah vs Tujuan", "status": "Valid" | "Warning" | "Critical", "feedback": "..." }},
                {{ "pair": "Metode vs Masalah", "status": "Valid" | "Warning" | "Critical", "feedback": "..." }},
                {{ "pair": "Tujuan vs Kesimpulan", "status": "Valid" | "Warning" | "Critical", "feedback": "..." }}
            ],
            "suggestions": ["Saran 1", "Saran 2"]
        }}
        """

        try:
            response = completion(
                model=model_name,
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.1,
                response_format={"type": "json_object"},
                max_tokens=2000
            )
            content = response.choices[0].message.content
            clean_json = ai_utils.clean_json_output(content)
            return json.loads(clean_json)

        except Exception as e:
            logger.error(f"Logic Check Error: {e}")
            return {"error": str(e)}

    # Legacy wrappers / Utils
    @staticmethod
    def paraphrase(user, text, style='academic', model='fast'):
        return ai_utils.paraphrase_text(user, text, style, model)

    @staticmethod
    def analyze_style_from_file(user, file_storage):
        try:
            return ai_utils.analyze_writing_style(file_storage.read().decode('utf-8', errors='ignore'), user.is_pro)
        except: return "Gagal analisis."

    @staticmethod
    def analyze_logic_flow(user, data):
        """
        Menganalisis konsistensi logis antara 3 elemen kunci:
        1. Judul (Topik Utama)
        2. Rumusan Masalah (Pertanyaan)
        3. Tujuan Penelitian (Jawaban Harapan)
        """
        title = data.get('title', '-')
        problem = data.get('problem', '-')
        objectives = data.get('objectives', '-')
        
        # Pilih Model Cerdas (Logic Check butuh nalar tinggi)
        # Gunakan 'logic_check' sebagai task_type agar router memilih GPT-4o (Pro) / Llama 70B (Free)
        model_name = get_smart_model('logic_check', user)

        system_instruction = (
            "Anda adalah Auditor Metodologi Penelitian Akademik yang sangat teliti.\n"
            "TUGAS: Periksa 'Benang Merah' (Konsistensi Logis) antara Judul, Rumusan Masalah, dan Tujuan.\n"
            "OUTPUT: JSON Valid berisi analisis per poin dan skor konsistensi."
        )

        user_prompt = f"""
        DATA SKRIPSI:
        1. JUDUL: "{title}"
        2. RUMUSAN MASALAH: "{problem}"
        3. TUJUAN PENELITIAN: "{objectives}"

        KRITERIA AUDIT:
        1. Apakah Rumusan Masalah menjawab variabel yang ada di Judul?
        2. Apakah Tujuan Penelitian sinkron/menjawab Rumusan Masalah? (Harus cerminan langsung).
        3. Apakah metode (tersirat dari judul) cocok dengan pertanyaan masalah?

        FORMAT OUTPUT JSON (STRICT):
        {{
            "consistency_score": 85, // Skala 0-100
            "status": "Cukup Konsisten", // Valid / Warning / Critical
            "analysis": [
                {{
                    "pair": "Judul vs Masalah",
                    "status": "Valid", // atau "Tidak Sinkron"
                    "feedback": "Penjelasan singkat..."
                }},
                {{
                    "pair": "Masalah vs Tujuan",
                    "status": "Warning",
                    "feedback": "Tujuan nomor 2 tidak menjawab masalah nomor 2..."
                }}
            ],
            "suggestions": [
                "Saran perbaikan konkret 1...",
                "Saran perbaikan konkret 2..."
            ]
        }}
        """

        try:
            response = completion(
                model=model_name,
                messages=[
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.2, # Rendah agar analisis stabil/objektif
                response_format={"type": "json_object"},
                drop_params=True
            )
            
            content = response.choices[0].message.content
            # Bersihkan JSON dari markdown wrapper ```json ... ```
            clean_content = clean_json_output(content)
            
            return json.loads(clean_content)

        except Exception as e:
            logger.error(f"Logic Audit AI Error: {e}")
            # Return fallback agar frontend tidak crash
            return {
                "consistency_score": 0,
                "status": "Error",
                "analysis": [],
                "suggestions": ["Gagal menganalisis logika. Silakan coba lagi."]
            }