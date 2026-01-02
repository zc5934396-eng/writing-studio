# File: app/utils/ppt_utils.py
# Deskripsi: Engine untuk mengubah Teks Skripsi menjadi File PowerPoint (.pptx)

import json
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from litellm import completion
from app.utils.ai_utils import AVAILABLE_MODELS, clean_json_output

def generate_slide_content(text_content, num_slides=10):
    """
    Menggunakan AI untuk meringkas teks skripsi menjadi struktur Slide Presentasi.
    Output: List of Dict [{'title': '...', 'points': ['...', '...']}]
    """
    model_name = AVAILABLE_MODELS['smart'] # Pakai model pintar biar ringkasannya tajam
    
    system_prompt = """
    PERAN: Anda adalah Pakar Presentasi Akademik.
    TUGAS: Ubah naskah skripsi/jurnal yang panjang menjadi outline presentasi yang padat, visual, dan menarik.
    
    FORMAT OUTPUT WAJIB (JSON ARRAY):
    [
        {
            "title": "Judul Slide (Singkat & Menarik)",
            "points": [
                "Poin utama 1 (maks 15 kata)",
                "Poin utama 2 (maks 15 kata)",
                "Poin data/angka penting"
            ]
        },
        ...
    ]
    
    ATURAN:
    1. Buat minimal 5-7 slide utama (Pendahuluan, Masalah, Metode, Hasil, Kesimpulan).
    2. Jangan copy-paste paragraf panjang! Gunakan bullet points yang tajam.
    3. Bahasa Indonesia baku tapi luwes untuk presentasi.
    """
    
    user_prompt = f"""
    [NASKAH SKRIPSI]
    {text_content[:15000]} 
    
    (Potongan teks di atas. Jika terpotong, simpulkan dari yang ada).
    
    Buatkan struktur presentasi sekitar {num_slides} slide.
    """

    try:
        response = completion(
            model=model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.3,
            response_format={"type": "json_object"}
        )
        
        content = clean_json_output(response.choices[0].message.content)
        parsed = json.loads(content)
        
        # Handle variasi output JSON
        if isinstance(parsed, dict):
            for key in ['slides', 'presentation', 'outline']:
                if key in parsed: return parsed[key]
            # Fallback values
            return list(parsed.values())[0] if parsed else []
            
        return parsed if isinstance(parsed, list) else []

    except Exception as e:
        print(f"PPT AI Error: {e}")
        return []

def create_presentation_file(slide_data, style='clean'):
    """
    Menggunakan python-pptx untuk membuat file .pptx dari data JSON.
    """
    prs = Presentation()
    
    # --- STYLE CONFIG ---
    # Warna default (Clean Blue/White)
    bg_color = RGBColor(255, 255, 255)
    title_color = RGBColor(0, 51, 102) # Dark Blue
    text_color = RGBColor(50, 50, 50)
    
    if style == 'dark':
        bg_color = RGBColor(30, 30, 35) # Dark Grey
        title_color = RGBColor(255, 200, 80) # Gold
        text_color = RGBColor(230, 230, 230) # White Smoke
    
    elif style == 'creative':
        bg_color = RGBColor(240, 248, 255) # Alice Blue
        title_color = RGBColor(108, 93, 211) # OnThesis Purple
        text_color = RGBColor(45, 48, 62)

    # 1. SLIDE JUDUL (Cover)
    if slide_data:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Apply Background (Manual simple fill)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        first_slide = slide_data[0]
        title.text = first_slide.get('title', 'Presentasi Sidang Skripsi')
        subtitle.text = "Dibuat Otomatis oleh OnThesis AI"
        
        # Styling Title
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True

    # 2. SLIDE ISI (Content)
    content_slide_layout = prs.slide_layouts[1] # Title + Content
    
    for slide_info in slide_data[1:]:
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Title
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get('title', 'Tanpa Judul')
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Bullet Points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        points = slide_info.get('points', [])
        # Jika poin berupa string panjang, pecah jadi list (jaga-jaga)
        if isinstance(points, str): points = [points]

        tf.clear() # Hapus placeholder default
        
        for point in points:
            p = tf.add_paragraph()
            p.text = str(point)
            p.font.size = Pt(20)
            p.font.color.rgb = text_color
            p.space_after = Pt(10)
            p.level = 0

    # 3. Simpan ke Memory Buffer
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io